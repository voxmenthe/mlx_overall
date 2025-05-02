# Copyright (c) Meta Platforms, Inc. and affiliates.
# All rights reserved.
#
# This source code is licensed under the terms described in the LICENSE file in
# the root directory of this source tree.
# PPTX parser logic

import os
from typing import Dict, Any

class PPTParser:
    """Parser for PowerPoint presentations"""
    
    def parse(self, file_path: str) -> str:
        """Parse a PPTX file into plain text
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Extracted text from the presentation
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX parsing. Install it with: pip install python-pptx")
        
        prs = Presentation(file_path)
        
        # Extract text from slides
        all_text = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- Slide {i+1} ---")
            
            # Get slide title
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Title: {slide.shapes.title.text}")
            
            # Get text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            all_text.append("\n".join(slide_text))
        
        return "\n\n".join(all_text)
    
    def save(self, content: str, output_path: str) -> None:
        """Save the extracted text to a file
        
        Args:
            content: Extracted text content
            output_path: Path to save the text
        """
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)